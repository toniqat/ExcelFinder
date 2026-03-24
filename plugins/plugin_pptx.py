"""PowerPoint 파일 플러그인 (python-pptx 필요) - Extension Pack"""
import sys
import pandas as pd
from pathlib import Path
from typing import List, Tuple

_src = Path(__file__).parent.parent / "src"
if str(_src) not in sys.path:
    sys.path.insert(0, str(_src))

from plugin_base import FormatPlugin


class PptxPlugin(FormatPlugin):
    """PowerPoint 파일(.pptx, .ppt) 처리 플러그인 (python-pptx 필요)

    .pptx: python-pptx로 슬라이드 텍스트 추출
    .ppt: 레거시 바이너리 형식 — python-pptx 미지원, 명확한 오류 반환
    """

    @property
    def plugin_id(self) -> str:
        return "pptx"

    @property
    def display_name(self) -> str:
        return "PowerPoint"

    @property
    def required_packages(self) -> List[str]:
        return ["pptx"]

    def supported_extensions(self) -> Tuple[str, ...]:
        return ('.pptx', '.ppt')

    def read_file(self, file_path: str) -> List[Tuple[str, List[str], pd.DataFrame]]:
        ext = Path(file_path).suffix.lower()

        if ext == '.ppt':
            raise ValueError(
                f".ppt (구형 바이너리 PowerPoint) 형식은 python-pptx로 읽을 수 없습니다: {file_path}\n"
                "파일을 .pptx 형식으로 변환한 후 다시 시도하세요."
            )

        from pptx import Presentation  # lazy import

        prs = Presentation(file_path)
        file_stem = Path(file_path).stem
        rows = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # 슬라이드 제목 추출 (없으면 'Slide N' 폴백)
            slide_title = f"Slide {slide_num}"
            try:
                title_shape = slide.shapes.title
                if title_shape and title_shape.has_text_frame:
                    t = title_shape.text.strip()
                    if t:
                        slide_title = t
            except Exception:
                pass

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        rows.append((str(slide_num), slide_title, text))

        if not rows:
            return []

        headers = ['Slide', 'Title', 'Content']
        df = pd.DataFrame(rows, columns=range(3))
        return [(file_stem, headers, df)]
